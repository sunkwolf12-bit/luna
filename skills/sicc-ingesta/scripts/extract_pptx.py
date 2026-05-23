#!/usr/bin/env python3
"""Wrapper documentado: extrae texto + imagenes de un PPTX a `/tmp/sicc/`.

El flujo preferido para Luna es `sicc parse --pptx <ruta>` (T2.5), que
hace lo mismo y ademas estructura el JSON candidato. Este script es solo
un helper standalone por si la skill necesita inspeccionar un PPTX sin
pasar por el CLI (debugging, troubleshoot, sesion sin CLI instalado).

Uso:

    python3 extract_pptx.py REPORTE_COBRANZA_JUNIO_2025.pptx \\
        [--out /tmp/sicc/<sesion>]

Produce:

- `/tmp/sicc/<sesion>/slides.json` con lista
  `{slide, titulo, total_apartado, imagenes: [ruta_png, ...]}`.
- `/tmp/sicc/<sesion>/slide_<n>_<slug>.png` por cada imagen embebida.

NO escribe a la DB. NO consolida. Solo extrae.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
import unicodedata
import uuid
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:  # pragma: no cover
    sys.stderr.write(
        "Falta python-pptx. Instalar con `pip install --user python-pptx`.\n"
    )
    sys.exit(2)


TOTAL_RE = re.compile(r"TOTAL\s*\$?\s*([\d,.]+)", re.IGNORECASE)


def _slug(s: str) -> str:
    norm = unicodedata.normalize("NFKD", s).encode("ascii", "ignore").decode()
    return re.sub(r"[^a-zA-Z0-9]+", "_", norm).strip("_").lower() or "slide"


def extract(pptx_path: Path, out_dir: Path) -> dict[str, Any]:
    out_dir.mkdir(parents=True, exist_ok=True)
    pres = Presentation(str(pptx_path))
    slides_out: list[dict[str, Any]] = []

    for idx, slide in enumerate(pres.slides, start=1):
        texts: list[str] = []
        imagenes: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blob = shape.image.blob
                ext = shape.image.ext or "png"
                titulo_tentativo = texts[0] if texts else f"slide{idx}"
                slug = _slug(titulo_tentativo)
                fname = f"slide_{idx:02d}_{slug}.{ext}"
                fpath = out_dir / fname
                fpath.write_bytes(blob)
                imagenes.append(str(fpath))

        titulo = texts[0] if texts else None
        total: str | None = None
        for t in texts:
            m = TOTAL_RE.search(t)
            if m:
                total = m.group(1).replace(",", "")
                break

        slides_out.append(
            {
                "slide": idx,
                "titulo": titulo,
                "total_apartado": total,
                "textos": texts,
                "imagenes": imagenes,
            }
        )

    result = {
        "pptx": str(pptx_path),
        "out_dir": str(out_dir),
        "slides": slides_out,
    }
    (out_dir / "slides.json").write_text(
        json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    return result


def main(argv: list[str]) -> int:
    ap = argparse.ArgumentParser(description=__doc__.split("\n\n", maxsplit=1)[0])
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--out", type=Path, default=None)
    args = ap.parse_args(argv)
    if not args.pptx.exists():
        sys.stderr.write(f"No existe {args.pptx}\n")
        return 1
    out_dir = args.out or Path(f"/tmp/sicc/{uuid.uuid4().hex[:8]}")
    result = extract(args.pptx, out_dir)
    json.dump(result, sys.stdout, ensure_ascii=False, indent=2)
    sys.stdout.write("\n")
    return 0


if __name__ == "__main__":  # pragma: no cover
    sys.exit(main(sys.argv[1:]))
